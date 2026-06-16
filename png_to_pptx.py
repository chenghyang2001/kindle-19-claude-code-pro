"""
PNG 投影片截圖合成 PPTX 工具

把 chapter-pptx/png/ 下的 PNG 截圖按章節分組，
合成真正的 PPTX 並覆蓋 chapter-pptx/ 下的同名 .pptx 檔案。
"""
import os
import re
import sys
import tempfile
from collections import defaultdict
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("錯誤：缺少 python-pptx，請執行 pip install python-pptx", file=sys.stderr)
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    print("錯誤：缺少 Pillow，請執行 pip install Pillow", file=sys.stderr)
    sys.exit(1)


# 16:9 寬螢幕投影片尺寸（EMU）：33.867cm × 19.05cm
SLIDE_WIDTH = Emu(12192000)
SLIDE_HEIGHT = Emu(6858000)

# 投影片序號正規式：檔名末尾的 -sNNN（一或多位數字）
_SLIDE_NUM_RE = re.compile(r"-s(\d+)$")


def parse_slide_info(png_path: Path):
    """
    從 PNG 路徑解析章節 stem 和投影片序號。

    命名格式：<chapter_stem>-s<NNN>.png
    範例：第01章_從聊天機器人到開發夥伴-s001.png
      → ('第01章_從聊天機器人到開發夥伴', 1)

    回傳 (chapter_stem, slide_num) 或 (None, None)（解析失敗）。
    為什麼取最後一個 -sNNN：附錄檔名如「附錄A_快速參考提示詞庫-s003.png」
    也符合格式，確保 stem 擷取完整章節名稱。
    """
    match = _SLIDE_NUM_RE.search(png_path.stem)
    if not match:
        return None, None
    slide_num = int(match.group(1))
    chapter_stem = png_path.stem[: match.start()]
    return chapter_stem, slide_num


def collect_png_groups(png_dir: Path) -> dict:
    """
    掃描 png_dir，按章節 stem 分組並依序號排序。

    回傳 {chapter_stem: [(slide_num, png_path), ...]}（各組已排序）。
    拋出 FileNotFoundError 若目錄不存在；ValueError 若沒有符合格式的 PNG。
    """
    if not png_dir.is_dir():
        raise FileNotFoundError(f"PNG 目錄不存在：{png_dir}")

    groups = defaultdict(list)
    skipped = 0

    for png_path in sorted(png_dir.glob("*.png")):
        chapter_stem, slide_num = parse_slide_info(png_path)
        if chapter_stem is None:
            print(f"  [略過] 無法解析序號：{png_path.name}", file=sys.stderr)
            skipped += 1
            continue

        # 用 PIL 驗證 PNG 完整性，損壞檔案跳過以免 PPTX 建立失敗
        try:
            with Image.open(png_path) as img:
                img.verify()
        except Exception as e:
            print(f"  [略過] PNG 損壞：{png_path.name} — {e}", file=sys.stderr)
            skipped += 1
            continue

        groups[chapter_stem].append((slide_num, png_path))

    if not groups:
        raise ValueError(
            f"PNG 目錄中找不到符合 '<章節>-s<NNN>.png' 格式的有效 PNG：{png_dir}"
        )

    # 按投影片序號升冪排序（確保 s001, s002... 順序正確）
    for stem in groups:
        groups[stem].sort(key=lambda x: x[0])

    if skipped:
        print(f"  （共略過 {skipped} 個不合格的 PNG）\n", file=sys.stderr)

    return dict(groups)


def get_blank_layout(prs: Presentation):
    """
    從 Presentation 取得空白版面，避免預設佔位符壓在 PNG 上方。
    優先選名稱含 'Blank' 或 '空白' 的；退路：index 6（Office 主題預設 Blank）。
    """
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower() or "空白" in layout.name:
            return layout
    # Office 主題通常 index 6 是 Blank，安全退路
    try:
        return prs.slide_layouts[6]
    except IndexError:
        return prs.slide_layouts[-1]


def build_chapter_pptx(
    chapter_stem: str,
    ordered_slides: list,
    output_dir: Path,
) -> Path:
    """
    為單一章節建立 PPTX 檔案。
    每張 PNG 對應一個 slide，圖片填滿整個投影片（left=0, top=0）。
    若 output_dir 下已有同名 .pptx，直接覆蓋（設計目標）。
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = get_blank_layout(prs)

    for _slide_num, png_path in ordered_slides:
        slide = prs.slides.add_slide(blank_layout)
        # 圖片填滿整個投影片（left=0, top=0）
        slide.shapes.add_picture(
            str(png_path),
            left=Emu(0),
            top=Emu(0),
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT,
        )

    output_path = output_dir / f"{chapter_stem}.pptx"

    # 原子性替換：先存暫存檔，成功後用 rename 替換目標
    # 避免磁碟滿或 Ctrl+C 中途中斷造成目標 .pptx 半寫損壞（ZIP 不完整）
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx", dir=output_path.parent)
    os.close(tmp_fd)
    try:
        prs.save(tmp_path)
        Path(tmp_path).replace(output_path)
    except Exception:
        Path(tmp_path).unlink(missing_ok=True)
        raise

    return output_path


def main() -> None:
    """
    主流程：
    1. 確認目錄結構
    2. 掃描 PNG 並按章節分組
    3. 逐章合成 PPTX（覆蓋既有 .pptx）
    4. 彙報成功與失敗數量
    """
    try:
        # 可攜路徑：不硬編碼使用者名稱，跨機器皆可執行
        base_dir = Path.home() / "workspace" / "kindle-19-claude-code-pro"
        png_dir = base_dir / "chapter-pptx" / "png"
        pptx_dir = base_dir / "chapter-pptx"

        print(f"PNG 來源目錄：{png_dir}")
        print(f"PPTX 輸出目錄：{pptx_dir}")
        print()

        groups = collect_png_groups(png_dir)
        print(f"偵測到 {len(groups)} 個章節，開始合成 PPTX...\n")

        success_count = 0
        error_count = 0

        for chapter_stem in sorted(groups.keys()):
            ordered_slides = groups[chapter_stem]
            try:
                output_path = build_chapter_pptx(chapter_stem, ordered_slides, pptx_dir)
                print(f"  完成：{output_path}（{len(ordered_slides)} 張投影片）")
                success_count += 1
            except Exception as e:
                print(f"  失敗：{chapter_stem} — {e}", file=sys.stderr)
                error_count += 1

        print(f"\n合成完成：{success_count} 個章節成功，{error_count} 個失敗。")

        if error_count > 0:
            sys.exit(1)

    except FileNotFoundError as e:
        print(f"錯誤：{e}", file=sys.stderr)
        sys.exit(1)
    except ValueError as e:
        print(f"錯誤：{e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"未預期的錯誤：{e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
